#!/usr/bin/env python3
"""Generates the Amidel eTender Scraper client-facing product presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1C, 0x38, 0x80)
NAVY_LIGHT  = RGBColor(0x25, 0x44, 0x99)
ORANGE      = RGBColor(0xF5, 0xA0, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF4, 0xF6, 0xFA)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x44)
MID_GRAY    = RGBColor(0x66, 0x66, 0x77)
LIGHT_GRAY  = RGBColor(0xE0, 0xE4, 0xEE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── Low-level helpers ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, font_name="Segoe UI",
                 italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, font_name="Segoe UI", italic=False,
             space_before_pt=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def navy_header_band(slide, title, subtitle=None):
    band_h = Inches(1.45)
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill_rgb=NAVY)
    add_rect(slide, 0, band_h - Inches(0.06), SLIDE_W, Inches(0.06), fill_rgb=ORANGE)
    add_text_box(slide, title,
                 left=Inches(0.55), top=Inches(0.18),
                 width=Inches(12.2), height=Inches(0.7),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     left=Inches(0.55), top=Inches(0.82),
                     width=Inches(12.2), height=Inches(0.4),
                     font_size=14, color=ORANGE, italic=True, align=PP_ALIGN.LEFT)


def content_area_top():
    return Inches(1.65)


def slide_footer(slide, page_num, total):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), fill_rgb=NAVY)
    add_text_box(slide, "Amidel (Pty) Ltd  |  Beyond Technology  |  Tender Intelligence",
                 left=Inches(0.4), top=SLIDE_H - Inches(0.31),
                 width=Inches(10), height=Inches(0.3),
                 font_size=9, color=OFF_WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, f"{page_num} / {total}",
                 left=Inches(12.5), top=SLIDE_H - Inches(0.31),
                 width=Inches(0.7), height=Inches(0.3),
                 font_size=9, color=ORANGE, align=PP_ALIGN.RIGHT)


def bullet_card(slide, left, top, width, height, heading, bullets,
                heading_size=15, bullet_size=12.5):
    add_rect(slide, left, top, width, height,
             fill_rgb=OFF_WHITE, line_rgb=LIGHT_GRAY, line_width_pt=0.75)
    add_rect(slide, left, top, Inches(0.055), height, fill_rgb=ORANGE)
    tx_left  = left + Inches(0.18)
    tx_width = width - Inches(0.22)
    txBox = slide.shapes.add_textbox(tx_left, top + Inches(0.14), tx_width, height - Inches(0.14))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = heading
    run.font.size = Pt(heading_size)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Segoe UI"
    for b in bullets:
        add_para(tf, "• " + b, font_size=bullet_size, color=DARK_GRAY)


# ── SLIDE BUILDERS ─────────────────────────────────────────────────────────────

TOTAL_SLIDES = 10


def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=NAVY)
    add_rect(slide, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5), fill_rgb=ORANGE)
    add_rect(slide, SLIDE_W - Inches(3.4), 0, Inches(3.4), SLIDE_H - Inches(0.5),
             fill_rgb=NAVY_LIGHT)
    add_rect(slide, SLIDE_W - Inches(3.4), 0, Inches(0.055), SLIDE_H - Inches(0.5),
             fill_rgb=ORANGE)

    # Product category chip
    add_rect(slide, Inches(0.7), Inches(1.7), Inches(3.2), Inches(0.4), fill_rgb=ORANGE)
    add_text_box(slide, "PROCUREMENT INTELLIGENCE  ·  SOUTH AFRICA",
                 left=Inches(0.72), top=Inches(1.72),
                 width=Inches(3.2), height=Inches(0.38),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, "Amidel\neTender Scraper",
                 left=Inches(0.7), top=Inches(2.2),
                 width=Inches(9.5), height=Inches(2.0),
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tagline
    add_text_box(slide, "Never miss a government tender again.",
                 left=Inches(0.7), top=Inches(4.25),
                 width=Inches(9), height=Inches(0.5),
                 font_size=20, bold=False, color=ORANGE, italic=True,
                 align=PP_ALIGN.LEFT)

    # Description
    add_text_box(slide,
                 "Automated daily monitoring of etenders.gov.za — structured reports, "
                 "department-level breakdowns, and actionable tender intelligence "
                 "delivered to your team without lifting a finger.",
                 left=Inches(0.7), top=Inches(4.88),
                 width=Inches(9.0), height=Inches(1.4),
                 font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Right panel
    add_text_box(slide, "Beyond\nTechnology",
                 left=SLIDE_W - Inches(3.2), top=Inches(2.8),
                 width=Inches(3.0), height=Inches(1.2),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Amidel (Pty) Ltd",
                 left=SLIDE_W - Inches(3.2), top=Inches(4.1),
                 width=Inches(3.0), height=Inches(0.4),
                 font_size=13, color=ORANGE, italic=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Confidential  ·  2026",
                 left=Inches(0.7), top=SLIDE_H - Inches(0.46),
                 width=Inches(9), height=Inches(0.4),
                 font_size=10, color=WHITE, align=PP_ALIGN.LEFT)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "The Challenge",
                     "Monitoring government tenders manually is unsustainable")
    slide_footer(slide, 2, TOTAL_SLIDES)

    top = content_area_top()
    pad   = Inches(0.45)
    col_w = Inches(5.9)
    gap   = Inches(0.55)

    problems = [
        ("Time-consuming",
         ["The eTenders portal must be checked every business day",
          "Tenders span multiple pages — no bulk export available",
          "Each listing must be opened individually to see full details"]),
        ("Easy to miss opportunities",
         ["New tenders appear without notification",
          "High-volume days mean relevant tenders get overlooked",
          "Closing dates are easy to miss when checking manually"]),
        ("No structured view",
         ["Website data is unstructured and inconsistent",
          "ICT vs RFQ classification requires manual judgement",
          "Briefing session details are buried inside each listing"]),
        ("No historical tracking",
         ["No built-in way to compare tender volumes across periods",
          "Department-level counts must be compiled by hand",
          "There is no audit trail of what was active and when"]),
    ]

    card_h = Inches(1.55)
    for i, (heading, bullets) in enumerate(problems):
        col = i % 2
        row = i // 2
        left = pad + col * (col_w + gap)
        t    = top + row * (card_h + Inches(0.18))
        bullet_card(slide, left, t, col_w, card_h, heading, bullets,
                    heading_size=14, bullet_size=12)


def slide_03_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "The Solution",
                     "One click. Structured output. Every time.")
    slide_footer(slide, 3, TOTAL_SLIDES)

    top = content_area_top()

    add_text_box(slide,
                 "The Amidel eTender Scraper automatically monitors etenders.gov.za on your behalf — "
                 "collecting every relevant tender, classifying it, and delivering clean, "
                 "structured reports directly to your team twice a week.",
                 left=Inches(0.7), top=top,
                 width=Inches(11.9), height=Inches(0.95),
                 font_size=15, color=DARK_GRAY, align=PP_ALIGN.LEFT)

    pill_top = top + Inches(1.1)
    pill_w   = Inches(3.6)
    pill_h   = Inches(3.8)
    gap      = Inches(0.42)

    pillars = [
        ("Monitor",
         "Automatically visits the eTenders portal for each day in the reporting period — "
         "no manual browsing required.",
         ["Runs across Mon–Wed and Thu–Sun periods",
          "Covers every tender listed, every day",
          "Handles large volumes without errors",
          "Duplicate detection built in"]),
        ("Classify",
         "Every tender is categorised by type (RFQ, ICT, Open Tender) and matched "
         "to tracked departments automatically.",
         ["26 organs of state tracked",
          "ICT and RFQ tenders flagged",
          "Briefing session details captured",
          "Department-level breakdowns produced"]),
        ("Deliver",
         "Structured Excel reports are generated and saved immediately after each run — "
         "ready to share with your team.",
         ["Per-department Tender Summary",
          "Combined RFQ & ICT Checker file",
          "Rolling 6-batch trend tracker",
          "Clickable links to source documents"]),
    ]

    for i, (title, desc, bullets) in enumerate(pillars):
        left = Inches(0.55) + i * (pill_w + gap)

        add_rect(slide, left, pill_top, pill_w, pill_h, fill_rgb=NAVY)
        add_rect(slide, left, pill_top, pill_w, Inches(0.07), fill_rgb=ORANGE)

        add_text_box(slide, title,
                     left=left + Inches(0.2), top=pill_top + Inches(0.15),
                     width=pill_w - Inches(0.3), height=Inches(0.5),
                     font_size=18, bold=True, color=WHITE)

        add_text_box(slide, desc,
                     left=left + Inches(0.2), top=pill_top + Inches(0.68),
                     width=pill_w - Inches(0.3), height=Inches(1.1),
                     font_size=11, color=LIGHT_GRAY)

        txBox = slide.shapes.add_textbox(
            left + Inches(0.2), pill_top + Inches(1.85),
            pill_w - Inches(0.3), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "✓  " + b
            run.font.size = Pt(11.5)
            run.font.color.rgb = ORANGE
            run.font.name = "Segoe UI"


def slide_04_use_cases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "Who Benefits",
                     "Any organisation that competes for — or tracks — government procurement")
    slide_footer(slide, 4, TOTAL_SLIDES)

    top = content_area_top()

    cases = [
        ("Sales & Business Development",
         ["Never miss an RFQ or open tender in your sector",
          "Get department-level counts delivered twice a week",
          "Act on opportunities before competitors notice them",
          "Briefing dates and compulsory attendance flagged upfront"]),
        ("Procurement Intelligence Teams",
         ["Track tender volumes across 26 organs of state",
          "Identify patterns — which departments are most active",
          "Compare ICT vs RFQ ratios across reporting periods",
          "Historical 6-batch rolling view for trend analysis"]),
        ("Executive Reporting",
         ["Clean, audit-ready Excel reports every batch cycle",
          "Summary counts ready to paste into board packs",
          "Consistent format — same layout every time",
          "No manual work between portal and report"]),
        ("Compliance & Supplier Development",
         ["Monitor tenders from specific municipalities and SOEs",
          "Track closing dates to ensure timely submissions",
          "Identify compulsory briefing sessions in advance",
          "Province-level visibility for targeted follow-up"]),
    ]

    card_h = Inches(1.65)
    col_w  = Inches(5.9)
    gap    = Inches(0.55)
    pad    = Inches(0.45)

    for i, (heading, bullets) in enumerate(cases):
        col = i % 2
        row = i // 2
        left = pad + col * (col_w + gap)
        t    = top + row * (card_h + Inches(0.15))
        bullet_card(slide, left, t, col_w, card_h, heading, bullets,
                    heading_size=14, bullet_size=12)


def slide_05_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "How It Works",
                     "From launch to ready-to-present output in minutes")
    slide_footer(slide, 5, TOTAL_SLIDES)

    steps = [
        ("1", "Select Period",
         "Choose the reporting period (Mon–Wed or Thu–Sun) "
         "using the calendar. The correct date range is set automatically."),
        ("2", "Run",
         "Click Run. The scraper visits etenders.gov.za "
         "for each day in the period and collects all tender listings."),
        ("3", "Extract",
         "Every tender's details are captured — description, "
         "department, type, closing date, briefing info, and more."),
        ("4", "Report",
         "Per-day files are saved, then combined into the "
         "RFQ & ICT Checker and the department Tender Summary."),
        ("5", "Track",
         "The rolling tracker is updated with the new period's "
         "counts and saved — ready to share immediately."),
    ]

    step_w = Inches(2.2)
    step_h = Inches(3.6)
    gap    = Inches(0.16)
    total  = len(steps) * step_w + (len(steps) - 1) * gap
    start  = (SLIDE_W - total) / 2
    top    = content_area_top() + Inches(0.3)

    for i, (num, title, desc) in enumerate(steps):
        left = start + i * (step_w + gap)

        add_rect(slide, left, top, step_w, step_h, fill_rgb=NAVY)
        add_rect(slide, left, top, step_w, Inches(0.06), fill_rgb=ORANGE)

        add_text_box(slide, num,
                     left=left, top=top + Inches(0.18),
                     width=step_w, height=Inches(0.8),
                     font_size=36, bold=True, color=ORANGE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     left=left + Inches(0.1), top=top + Inches(1.0),
                     width=step_w - Inches(0.2), height=Inches(0.55),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, desc,
                     left=left + Inches(0.12), top=top + Inches(1.6),
                     width=step_w - Inches(0.24), height=Inches(1.8),
                     font_size=11, color=LIGHT_GRAY,
                     align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow_left = left + step_w + gap * 0.1
            arrow_top  = top + step_h / 2 - Inches(0.15)
            add_text_box(slide, "▶",
                         left=arrow_left, top=arrow_top,
                         width=gap * 0.8, height=Inches(0.3),
                         font_size=11, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_06_outputs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "What You Receive",
                     "Structured, consistent reports delivered after every run")
    slide_footer(slide, 6, TOTAL_SLIDES)

    top = content_area_top()

    outputs = [
        ("RFQ & ICT Checker",
         "The primary report — all tenders for the period in one workbook. "
         "Includes a full data tab and a summary tab with counts broken down "
         "by department, ICT category, and tender type."),
        ("Tender Summary",
         "One tab per tracked department that had tenders in the period. "
         "RFQs are sorted to the top of each tab, followed by all other "
         "tender types ordered by closing date. Includes briefing date and "
         "whether attendance is compulsory."),
        ("Per-Day Batch Files",
         "Individual Excel files for each day in the reporting period. "
         "Useful for auditing or drilling into a specific day's activity "
         "without opening the combined file."),
        ("Rolling Tracker",
         "A running comparison of NNT, ICT, and RFQ counts across the last "
         "six reporting periods — one column per batch, newest on the left. "
         "Automatically maintained and shared via OneDrive."),
        ("Clickable Source Links",
         "Every tender row contains a hyperlink directly to the official "
         "listing on etenders.gov.za — one click to access the full "
         "tender document and submission instructions."),
    ]

    card_h = Inches(1.0)
    card_w = Inches(12.3)
    gap    = Inches(0.12)
    pad    = Inches(0.5)

    for i, (title, desc) in enumerate(outputs):
        t = top + i * (card_h + gap)
        add_rect(slide, pad, t, card_w, card_h,
                 fill_rgb=OFF_WHITE, line_rgb=LIGHT_GRAY, line_width_pt=0.5)
        add_rect(slide, pad, t, Inches(0.05), card_h, fill_rgb=ORANGE)

        # Number badge
        add_rect(slide, pad + Inches(0.15), t + Inches(0.29),
                 Inches(0.36), Inches(0.36), fill_rgb=NAVY)
        add_text_box(slide, str(i + 1),
                     left=pad + Inches(0.15), top=t + Inches(0.28),
                     width=Inches(0.36), height=Inches(0.36),
                     font_size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     left=pad + Inches(0.65), top=t + Inches(0.08),
                     width=Inches(2.4), height=Inches(0.35),
                     font_size=12, bold=True, color=NAVY)
        add_text_box(slide, desc,
                     left=pad + Inches(0.65), top=t + Inches(0.42),
                     width=card_w - Inches(0.75), height=Inches(0.5),
                     font_size=10.5, color=DARK_GRAY)


def slide_07_data_fields(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "Data Captured",
                     "Every tender record contains up to 23 fields of structured information")
    slide_footer(slide, 7, TOTAL_SLIDES)

    top = content_area_top()

    fields = [
        ("Report Date",             "The reporting period this tender belongs to"),
        ("Tender ID",               "Official government tender reference number"),
        ("Publication Date",        "When the tender was published on eTenders"),
        ("Closing Date & Time",     "Submission deadline — date and exact time"),
        ("Tender Type",             "e.g. RFQ, Open Tender, Negotiation, etc."),
        ("Tender Description",      "Full title and description of the tender"),
        ("Department",              "Organ of state advertising the tender"),
        ("Province",                "Province the department falls under"),
        ("Category",                "Sector category (e.g. ICT, Construction)"),
        ("eSubmission",             "Whether electronic submission is available"),
        ("Briefing Session",        "Whether a briefing session has been scheduled"),
        ("Briefing Date",           "Date of the briefing session (Tender Summary only)"),
        ("Compulsory Briefing",     "Whether attendance at briefing is mandatory"),
        ("Briefing Venue",          "Location or details of the briefing session"),
        ("Link",                    "Direct clickable link to the tender document"),
        ("SOE Flag",                "Whether the entity is a state-owned enterprise"),
        ("Cost of Sales Estimate",  "Estimated value or cost attached to the tender"),
        ("Capability Available",    "Internal capability classification"),
        ("Requirements",            "Any special requirements listed on the tender"),
    ]

    col_w  = Inches(5.7)
    row_h  = Inches(0.3)
    gap    = Inches(0.55)
    pad    = Inches(0.45)

    half = (len(fields) + 1) // 2
    for i, (field, desc) in enumerate(fields):
        col = i // half
        row = i % half
        left = pad + col * (col_w + gap)
        t    = top + row * row_h

        add_rect(slide, left, t + Inches(0.03), Inches(2.5), Inches(0.24), fill_rgb=NAVY)
        add_text_box(slide, field,
                     left=left + Inches(0.05), top=t + Inches(0.03),
                     width=Inches(2.4), height=Inches(0.24),
                     font_size=9, bold=True, color=WHITE,
                     font_name="Segoe UI", align=PP_ALIGN.LEFT)
        add_text_box(slide, desc,
                     left=left + Inches(2.6), top=t + Inches(0.03),
                     width=col_w - Inches(2.65), height=Inches(0.24),
                     font_size=9.5, color=DARK_GRAY)


def slide_08_departments(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "Coverage",
                     "26 organs of state monitored every reporting period")
    slide_footer(slide, 8, TOTAL_SLIDES)

    top = content_area_top()

    depts = [
        "eTenders (All)", "Eskom", "SITA", "Transnet", "JPC",
        "Matatiele LM", "Ntabankulu LM", "Umzimvubu LM",
        "Winnie Mandela LM", "Mnquma LM", "Great Kei LM",
        "Amahlathi LM", "Raymond Mhlaba LM", "JOSHCO",
        "GPL", "RAF", "MM Trading Company",
        "Nelson Mandela Bay MM", "Buffalo City MM",
        "EC DPW", "DEL", "SIU", "GDoH", "DIRCO", "CP JHB", "W JHB",
    ]

    add_text_box(slide,
                 "The following departments and organs of state are tracked each batch. "
                 "For each one, the reports show the number of new tenders (NNT), "
                 "ICT tenders, and RFQs. The list can be customised to match your pipeline.",
                 left=Inches(0.5), top=top,
                 width=Inches(12.3), height=Inches(0.65),
                 font_size=13, color=DARK_GRAY)

    chip_w    = Inches(2.45)
    chip_h    = Inches(0.42)
    gap_x     = Inches(0.18)
    gap_y     = Inches(0.14)
    cols      = 5
    pad       = Inches(0.5)
    start_top = top + Inches(0.8)

    for i, dept in enumerate(depts):
        col  = i % cols
        row  = i // cols
        left = pad + col * (chip_w + gap_x)
        t    = start_top + row * (chip_h + gap_y)

        fill    = NAVY if i == 0 else OFF_WHITE
        txt_col = WHITE if i == 0 else NAVY

        add_rect(slide, left, t, chip_w, chip_h,
                 fill_rgb=fill, line_rgb=NAVY_LIGHT, line_width_pt=0.6)
        if i > 0:
            add_rect(slide, left, t, Inches(0.04), chip_h, fill_rgb=ORANGE)

        add_text_box(slide, dept,
                     left=left + (Inches(0.08) if i > 0 else 0),
                     top=t + Inches(0.07),
                     width=chip_w - Inches(0.08), height=chip_h - Inches(0.08),
                     font_size=11, bold=(i == 0), color=txt_col,
                     align=PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT)


def slide_09_why_amidel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_header_band(slide, "Why Amidel",
                     "Local expertise. Custom solutions. Real results.")
    slide_footer(slide, 9, TOTAL_SLIDES)

    top = content_area_top()

    reasons = [
        ("Built for South Africa",
         ["Designed specifically for etenders.gov.za",
          "Understands the SA government procurement landscape",
          "Tracks the organs of state relevant to your sector",
          "Adapts to local tender formats and classifications"]),
        ("Delivered, Not Just Built",
         ["Reports are generated and distributed automatically",
          "Your team receives structured outputs — no setup required",
          "OneDrive integration keeps everyone in sync",
          "Consistent format every batch, without fail"]),
        ("Customisable to Your Needs",
         ["Department watchlist tailored to your pipeline",
          "Additional classification fields can be added",
          "Reporting cadence matches your business rhythm",
          "Future modules: email delivery, dashboard views"]),
        ("Proven & Actively Maintained",
         ["Used internally at Amidel for live procurement tracking",
          "Refined through real-world batch cycles",
          "Robust retry logic handles network and portal issues",
          "Continuously improved based on team feedback"]),
    ]

    card_h = Inches(1.65)
    col_w  = Inches(5.9)
    gap    = Inches(0.55)
    pad    = Inches(0.45)

    for i, (heading, bullets) in enumerate(reasons):
        col = i % 2
        row = i // 2
        left = pad + col * (col_w + gap)
        t    = top + row * (card_h + Inches(0.15))
        bullet_card(slide, left, t, col_w, card_h, heading, bullets,
                    heading_size=14, bullet_size=12)


def slide_10_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=NAVY)
    add_rect(slide, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5), fill_rgb=ORANGE)
    add_rect(slide, SLIDE_W - Inches(3.4), 0, Inches(3.4), SLIDE_H - Inches(0.5),
             fill_rgb=NAVY_LIGHT)
    add_rect(slide, SLIDE_W - Inches(3.4), 0, Inches(0.055), SLIDE_H - Inches(0.5),
             fill_rgb=ORANGE)

    add_text_box(slide, "Let's Get Started",
                 left=Inches(0.7), top=Inches(0.9),
                 width=Inches(8.5), height=Inches(0.7),
                 font_size=40, bold=True, color=WHITE)
    add_text_box(slide, "What the Amidel eTender Scraper delivers for your organisation",
                 left=Inches(0.7), top=Inches(1.6),
                 width=Inches(8.5), height=Inches(0.4),
                 font_size=15, color=ORANGE, italic=True)

    points = [
        "Automated monitoring of 26 government departments — twice a week",
        "Clean, structured Excel reports ready to act on immediately",
        "Department-level tender counts with ICT and RFQ breakdowns",
        "Briefing session dates and compulsory attendance flagged upfront",
        "Rolling 6-period trend tracker — always up to date",
        "Customisable to your department watchlist and reporting needs",
    ]

    for i, pt in enumerate(points):
        add_text_box(slide, "✓  " + pt,
                     left=Inches(0.85), top=Inches(2.25) + i * Inches(0.62),
                     width=Inches(8.5), height=Inches(0.55),
                     font_size=14, color=WHITE)

    # Right panel — contact
    add_text_box(slide, "Contact\nAmidel",
                 left=SLIDE_W - Inches(3.2), top=Inches(1.4),
                 width=Inches(3.0), height=Inches(1.1),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    contact_lines = [
        "amidel.co.za",
        " ",
        "Request a demo",
        "or tailored setup",
        " ",
        "Beyond Technology",
    ]
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(3.15), Inches(2.65),
        Inches(2.9), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(contact_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12 if j not in (2, 3) else 11)
        run.font.color.rgb = ORANGE if j == 0 else (WHITE if j != 5 else LIGHT_GRAY)
        run.font.name = "Segoe UI"
        run.font.bold = j == 0
        run.font.italic = j == 5

    add_text_box(slide, "Amidel (Pty) Ltd  ·  2026",
                 left=Inches(0.7), top=SLIDE_H - Inches(0.46),
                 width=Inches(9), height=Inches(0.4),
                 font_size=10, color=WHITE)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_use_cases(prs)
    slide_05_workflow(prs)
    slide_06_outputs(prs)
    slide_07_data_fields(prs)
    slide_08_departments(prs)
    slide_09_why_amidel(prs)
    slide_10_cta(prs)

    out = "Amidel_eTender_Scraper_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
